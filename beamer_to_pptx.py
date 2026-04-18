#!/usr/bin/env python3
"""Convert a Beamer .tex presentation to an editable PPTX.

Usage:
    python3 beamer_to_pptx.py presentation.tex [presentation.pdf] [output.pptx]

The .pdf argument is optional but enables high-fidelity fallback rendering
for slides containing TikZ, tables, or other complex LaTeX that cannot be
directly translated to PowerPoint objects. When a slide uses these elements,
the corresponding PDF page is clip-rasterized around the element's own
vector/image bounding box so that surrounding captions or bullet text are
not pulled into the extracted image.

Figure collection
-----------------
The expected workflow is: the user drops the paper into ``Paper_files/``,
the Cursor agent extracts the figures it wants into ``Extracted_figures/``
(as part of slide generation), and this converter then runs on
``presentation.tex``. As a safety net on top of the agent's extraction,
every ``\\includegraphics`` path referenced in the .tex is resolved against
the ``\\graphicspath`` entries in the .tex (typically
``{Extracted_figures/}{Paper_files/}{Assets/}``) and **copied into**
``Extracted_figures/`` if it isn't already there. After conversion,
``Extracted_figures/`` holds a canonical copy of every figure used.
Auto-cropped TikZ/table PNGs are written as ``slide_content_<N>.png``.
The footer logo (``Assets/logo.jpg``) is *not* copied - it's chrome, not
slide content.

Speaker notes (``\\note{...}`` after a frame) and a footer logo are enabled
by default. The logo is auto-detected from the .tex (any ``\\includegraphics``
whose path contains ``logo``) and, failing that, from ``Assets/logo.jpg``
or ``Assets/logo.png``.

Style settings (font, colors, sizes, footer) are configured via constants
at the top of the file.
"""

import re
import os
import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# =====================================================================
# STYLE CONFIGURATION
# =====================================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.4)
CONTENT_L = MARGIN
CONTENT_R = SLIDE_W - MARGIN
CONTENT_W = CONTENT_R - CONTENT_L

TITLE_TOP = Inches(0.20)
TITLE_H = Inches(0.65)
BODY_TOP = Inches(1.35)
FOOTER_Y = SLIDE_H - Inches(0.48)
BODY_BOTTOM = FOOTER_Y - Inches(0.08)
BODY_H = BODY_BOTTOM - BODY_TOP

FONT_SLIDE_TITLE = Pt(36)
FONT_BODY = Pt(24)
FONT_SUB = Pt(20)
FONT_CAPTION = Pt(14)
FONT_FOOTER = Pt(12)

INDENT_L1 = Inches(0.35)
INDENT_L2 = Inches(0.80)

# Use a font installed on your system (PowerPoint falls back if missing)
FONT_NAME = "Calibri"
# Center footer when \\date{...} is missing from the .tex preamble
FOOTER_TEXT_FALLBACK = ""

BLACK  = RGBColor(0x00, 0x00, 0x00)
GRAY   = RGBColor(0x80, 0x80, 0x80)
ACCENT = RGBColor(0xFF, 0xC0, 0x00)

COL_GAP = Inches(0.25)
COL_SPLIT = 0.45

PDF_DPI = 300

# =====================================================================
# LATEX PARSER
# =====================================================================

def read_tex(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()


def extract_preamble(tex):
    """Pull metadata from preamble: title, author, institute, date, graphicspath."""
    info = {}
    m = re.search(r'\\title\{(.+?)\}', tex, re.DOTALL)
    if m:
        info['title'] = clean_text(m.group(1))
    m = re.search(r'\\author\{(.+?)\}(?=\s*\\institute|\s*\\date)', tex, re.DOTALL)
    if m:
        info['author'] = clean_text(m.group(1))
    m = re.search(r'\\institute\{(.+?)\}', tex, re.DOTALL)
    if m:
        info['institute'] = clean_text(m.group(1))
    m = re.search(r'\\date\{(.+?)\}', tex, re.DOTALL)
    if m:
        info['date'] = m.group(1).strip()
    m = re.search(r'\\graphicspath\s*\{((?:\s*\{[^{}]+\}\s*)+)\}', tex)
    if m:
        info['graphicspath'] = re.findall(r'\{([^{}]+)\}', m.group(1))
    return info


def clean_text(s):
    """Strip LaTeX formatting commands, keeping readable text."""
    s = re.sub(r'\\\\(\[.*?\])?', '\n', s)
    s = re.sub(r'\\begin\{(center|flushleft|flushright|block|alertblock|exampleblock)\}', '', s)
    s = re.sub(r'\\end\{(center|flushleft|flushright|block|alertblock|exampleblock)\}', '', s)
    s = re.sub(r'\\textbf\{(.*?)\}', r'**\1**', s)
    s = re.sub(r'\\textit\{(.*?)\}', r'_\1_', s)
    s = re.sub(r'\{\\(small|footnotesize|large|Large|LARGE|normalsize|tiny)\s*', '', s)
    s = re.sub(r'\\(small|footnotesize|large|Large|LARGE|normalsize|tiny)\b\s*', '', s)
    s = re.sub(r'\\textcolor\{.*?\}\{(.*?)\}', r'\1', s)
    s = re.sub(r'\\(par|centering|vfill|noindent)\b', '', s)
    s = re.sub(r'\\vspace\{.*?\}', '', s)
    s = re.sub(r'\\hspace\{.*?\}', '', s)
    s = re.sub(r'\\[, ;!]', ' ', s)
    s = s.replace('\\&', '&')
    s = s.replace('\\%', '%')
    s = s.replace('\\$', '$')
    s = s.replace('~', ' ')
    s = s.replace('``', '\u201c').replace("''", '\u201d')
    s = re.sub(r'\$([^$]+)\$', r'\1', s)
    s = re.sub(r'\\(?:geq|ge)\b', '\u2265', s)
    s = re.sub(r'\\(?:leq|le)\b', '\u2264', s)
    s = re.sub(r'\\(?:to|rightarrow)\b', '\u2192', s)
    s = re.sub(r'\\(?:times)\b', '\u00d7', s)
    s = re.sub(r'\\(?:uparrow)\b', '\u2191', s)
    s = re.sub(r'\\(?:downarrow)\b', '\u2193', s)
    s = re.sub(r'\\(?:neq|ne)\b', '\u2260', s)
    s = re.sub(r'\\(?:approx)\b', '\u2248', s)
    s = re.sub(r'\\(?:infty)\b', '\u221e', s)
    s = re.sub(r'\\[a-zA-Z]+\{(.*?)\}', r'\1', s)
    s = re.sub(r'\\[a-zA-Z]+', '', s)
    s = re.sub(r'[{}]', '', s)
    s = re.sub(r'[ \t]+', ' ', s)
    return s.strip()


def extract_frames(tex):
    """Return a list of frame dicts from the document body."""
    body_m = re.search(r'\\begin\{document\}(.*?)\\end\{document\}', tex, re.DOTALL)
    if not body_m:
        return []
    body = body_m.group(1)

    frames = []
    chunks = re.split(r'(?=\\begin\{frame\}|\\frame\{\\titlepage\})', body)
    note_pattern = re.compile(r'\\note\{((?:[^{}]|\{(?:[^{}]|\{[^{}]*\})*\})*)\}', re.DOTALL)

    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue

        frame = {}

        if '\\frame{\\titlepage}' in chunk:
            frame['type'] = 'titlepage'
            nm = note_pattern.search(chunk)
            if nm:
                frame['note'] = clean_text(nm.group(1))
            frames.append(frame)
            continue

        fm = re.search(r'\\begin\{frame\}(?:\{(.*?)\})?(.+?)\\end\{frame\}', chunk, re.DOTALL)
        if not fm:
            continue

        title = clean_text(fm.group(1)) if fm.group(1) else ''
        content = fm.group(2)

        frame['type'] = 'content'
        frame['title'] = title
        frame['raw'] = content

        has_tikz = '\\begin{tikzpicture}' in content
        has_table = '\\begin{tabular}' in content
        frame['needs_fallback'] = has_tikz or has_table

        has_columns = '\\begin{columns}' in content
        frame['has_columns'] = has_columns

        if has_columns:
            cols = re.findall(
                r'\\begin\{column\}\{.*?\}(.*?)\\end\{column\}',
                content, re.DOTALL
            )
            if len(cols) >= 2:
                frame['left_col'] = cols[0]
                frame['right_col'] = cols[1]

        frame['images'] = re.findall(
            r'\\includegraphics(?:\[.*?\])?\{(.*?)\}', content
        )
        frame['captions'] = re.findall(
            r'\\textit\{(.*?)\}', content
        )

        frame['items'] = parse_items(content)
        frame['is_enumerate'] = '\\begin{enumerate}' in content

        quote_m = re.search(r'\{\\large\s*\\textit\{(.+?)\}\}', content, re.DOTALL)
        if not quote_m:
            quote_m = re.search(r'\{\\LARGE\s*\\textbf\{(.+?)\}\}', content, re.DOTALL)
        if quote_m:
            frame['quote'] = clean_text(quote_m.group(1))

        rest = chunk[fm.end():]
        nm = note_pattern.search(rest)
        if not nm:
            nm = note_pattern.search(chunk)
        if nm:
            frame['note'] = clean_text(nm.group(1))

        frames.append(frame)

    return frames


def parse_items(content):
    """Parse nested itemize/enumerate into flat list of (text, level) tuples."""
    items = []
    level = 0
    in_list = False

    for line in content.split('\n'):
        stripped = line.strip()
        if re.match(r'\\begin\{(itemize|enumerate)\}', stripped):
            if in_list:
                level += 1
            in_list = True
            continue
        if re.match(r'\\end\{(itemize|enumerate)\}', stripped):
            level = max(0, level - 1)
            if level == 0:
                in_list = False
            continue
        if stripped.startswith('\\item'):
            text = stripped[5:].strip()
            text = clean_text(text)
            if text:
                items.append((text, level))

    return items


# =====================================================================
# PPTX BUILDER HELPERS (from gen_pptx.py)
# =====================================================================

def _zero_insets(tf):
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0


def _clear_bullet_tags(pPr):
    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum', 'a:buClr', 'a:buSzPct', 'a:buFont'):
        for old in pPr.findall(qn(tag)):
            pPr.remove(old)


def _set_bullet_common(paragraph, color, level, margin, indent):
    pPr = paragraph._pPr
    _clear_bullet_tags(pPr)
    pPr.set('marL', str(int(margin)))
    pPr.set('indent', str(-int(indent)))
    pPr.set('lvl', str(level))

    buClr = etree.SubElement(pPr, qn('a:buClr'))
    srgb = etree.SubElement(buClr, qn('a:srgbClr'))
    srgb.set('val', str(color))

    buSz = etree.SubElement(pPr, qn('a:buSzPct'))
    buSz.set('val', '120000')

    buFont = etree.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', FONT_NAME)
    return pPr


def _set_bullet_char(paragraph, char, color, level, margin, indent):
    pPr = _set_bullet_common(paragraph, color, level, margin, indent)
    el = etree.SubElement(pPr, qn('a:buChar'))
    el.set('char', char)


def _set_bullet_number(paragraph, color, level, margin, indent, num_type='arabicPeriod'):
    pPr = _set_bullet_common(paragraph, color, level, margin, indent)
    el = etree.SubElement(pPr, qn('a:buAutoNum'))
    el.set('type', num_type)


def add_footer(slide, slide_num, total_slides, logo_path, center_text=None):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), FOOTER_Y, SLIDE_W, Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()
    line.shadow.inherit = False

    ty = FOOTER_Y + Inches(0.06)
    th = Inches(0.32)

    if slide_num > 1:
        nb = slide.shapes.add_textbox(MARGIN, ty, Inches(1.5), th)
        _zero_insets(nb.text_frame)
        p = nb.text_frame.paragraphs[0]
        p.text = f"({slide_num - 1}/{total_slides - 1})"
        p.font.size = FONT_FOOTER; p.font.color.rgb = GRAY

    db = slide.shapes.add_textbox(SLIDE_W // 2 - Inches(1.5), ty, Inches(3), th)
    _zero_insets(db.text_frame)
    p = db.text_frame.paragraphs[0]
    p.text = (center_text if center_text is not None else FOOTER_TEXT_FALLBACK) or ""
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_FOOTER; p.font.color.rgb = GRAY

    if logo_path and os.path.exists(logo_path):
        from PIL import Image as PILImage
        with PILImage.open(logo_path) as img:
            aspect = img.width / img.height
        lh = Inches(0.24)
        lw = Inches(0.24 * aspect)
        slide.shapes.add_picture(logo_path, CONTENT_R - lw, ty + Inches(0.02),
                                 width=lw, height=lh)


def add_title_bar(slide, text):
    tb = slide.shapes.add_textbox(CONTENT_L, TITLE_TOP, CONTENT_W, TITLE_H)
    tf = tb.text_frame; tf.word_wrap = True; _zero_insets(tf)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = FONT_SLIDE_TITLE; p.font.bold = True; p.font.color.rgb = BLACK


def _add_runs_from_text(paragraph, text, font_size):
    """Parse **bold** and _italic_ markers into runs."""
    parts = re.split(r'(\*\*.*?\*\*|_.*?_)', text)
    for part in parts:
        if not part:
            continue
        r = paragraph.add_run()
        if part.startswith('**') and part.endswith('**'):
            r.text = part[2:-2]
            r.font.bold = True
        elif part.startswith('_') and part.endswith('_'):
            r.text = part[1:-1]
            r.font.italic = True
        else:
            r.text = part
        r.font.size = font_size
        r.font.color.rgb = BLACK


def add_bullet_list(slide, items, left, top, width, height, is_enum=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True; _zero_insets(tf)

    ENUM_L1 = Inches(0.60)
    ENUM_L2 = Inches(1.05)

    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)

        if is_enum and level == 0:
            font_sz = FONT_BODY
            _set_bullet_number(p, ACCENT, level, ENUM_L1, ENUM_L1)
        elif is_enum and level > 0:
            font_sz = FONT_SUB
            _set_bullet_char(p, '\u2013', ACCENT, level, ENUM_L2, INDENT_L1)
        elif level == 0:
            font_sz = FONT_BODY
            _set_bullet_char(p, '\u2022', ACCENT, level, INDENT_L1, INDENT_L1)
        else:
            font_sz = FONT_SUB
            _set_bullet_char(p, '\u2013', ACCENT, level, INDENT_L2, INDENT_L1)

        _add_runs_from_text(p, text, font_sz)


def add_quote(slide, text, top):
    tb = slide.shapes.add_textbox(CONTENT_L + Inches(0.3), top, CONTENT_W - Inches(0.6), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True; _zero_insets(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(22); r.font.italic = True; r.font.color.rgb = BLACK


def set_notes(slide, text):
    """Always attach a notes slide so the Notes pane is visible in PowerPoint.

    If *text* is empty we still create the notes slide (with empty text),
    which makes speaker notes a first-class default rather than an opt-in.
    """
    slide.notes_slide.notes_text_frame.text = text or ""


def apply_font_to_all(prs, font_name):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name


def resolve_image(name, graphicspath, base_dir):
    """Find an image file, trying each graphicspath entry and common extensions.

    *graphicspath* may be a list of path strings (from ``\\graphicspath{{a/}{b/}}``),
    a single string, or falsy. All entries are searched in order, followed by
    *base_dir* as a last resort.
    """
    candidates = [name]
    if not os.path.splitext(name)[1]:
        for ext in ['.png', '.jpg', '.jpeg', '.pdf']:
            candidates.append(name + ext)

    if isinstance(graphicspath, str):
        gp_list = [graphicspath] if graphicspath else []
    elif graphicspath:
        gp_list = list(graphicspath)
    else:
        gp_list = []

    dirs = [os.path.join(base_dir, p) for p in gp_list]
    dirs.append(base_dir)

    for d in dirs:
        for c in candidates:
            p = os.path.join(d, c)
            if os.path.exists(p):
                return p
    return None


def ensure_in_extracted(src_path, extracted_dir):
    """Copy *src_path* into ``Extracted_figures/`` if it isn't already there.

    Returns the path inside *extracted_dir* (or *src_path* unchanged when the
    source is already inside that folder or when the copy is a no-op). Files
    with identical size at the destination are left untouched so repeat runs
    don't churn timestamps.
    """
    import shutil

    if not src_path or not os.path.exists(src_path):
        return src_path
    os.makedirs(extracted_dir, exist_ok=True)
    src_abs = os.path.abspath(src_path)
    ef_abs = os.path.abspath(extracted_dir)
    try:
        if os.path.commonpath([src_abs, ef_abs]) == ef_abs:
            return src_path
    except ValueError:
        pass
    dst = os.path.join(extracted_dir, os.path.basename(src_path))
    if os.path.exists(dst):
        try:
            if os.path.getsize(dst) == os.path.getsize(src_path):
                return dst
        except OSError:
            pass
    shutil.copy2(src_path, dst)
    return dst


# =====================================================================
# PDF FALLBACK RENDERER
# =====================================================================

def rasterize_content_region(pdf_path, page_num, out_dir, has_bullets=False,
                             has_table=False, dpi=PDF_DPI):
    """Render the complex element (table / TikZ / figure) from a PDF page.

    Strategy
    --------
    Use the PDF's native primitives to locate the visual element directly:

    * drawings (vector paths) - captures booktabs rule lines and TikZ shapes
    * embedded image rects    - captures includegraphics figures

    The union of these rectangles, restricted to the body zone (below the
    frame title, above the footer), gives a tight bbox.  Text blocks that
    lie within that bbox's vertical range and horizontally overlap are
    merged in so that cell text, axis labels, and in-diagram annotations
    are preserved.

    When *has_bullets* is True (bullet list coexists with the diagram),
    text that sits below the bottom rule line or below the drawings' y1
    is assumed to be the bullet list and is excluded from the bbox.

    If no usable primitives are found the function falls back to a
    whitespace-band analysis of a full-page raster.
    """
    doc = fitz.open(pdf_path)
    if page_num >= len(doc):
        doc.close()
        return None

    page = doc[page_num]
    page_w, page_h = page.rect.width, page.rect.height
    title_h = page_h * 0.14
    footer_y = page_h * 0.90
    body_zone = fitz.Rect(0, title_h, page_w, footer_y)

    def _in_body(r):
        """Return rect clipped to body zone, or None if fully outside.

        Works for degenerate rects (zero-height lines) which fitz.intersects
        rejects.
        """
        x0 = max(r.x0, body_zone.x0)
        y0 = max(r.y0, body_zone.y0)
        x1 = min(r.x1, body_zone.x1)
        y1 = min(r.y1, body_zone.y1)
        if x1 < x0 or y1 < y0:
            return None
        return fitz.Rect(x0, y0, x1, y1)

    drawing_rects = []
    for d in page.get_drawings():
        dtype = d.get('type')
        if dtype == 'f':
            dr = d.get('rect')
            if dr and dr.width >= page_w * 0.9 and dr.height >= page_h * 0.9:
                continue
        items = d.get('items') or []
        if not items:
            r = d.get('rect')
            if r is None or r.is_empty:
                continue
            rr = _in_body(r)
            if rr is None:
                continue
            if rr.width >= page_w * 0.9 and rr.height >= body_zone.height * 0.9:
                continue
            if rr.width >= 1 or rr.height >= 1:
                drawing_rects.append(rr)
            continue
        for item in items:
            op = item[0]
            prim = None
            if op == 'l':
                p1, p2 = item[1], item[2]
                prim = fitz.Rect(
                    min(p1.x, p2.x), min(p1.y, p2.y),
                    max(p1.x, p2.x), max(p1.y, p2.y),
                )
            elif op == 're':
                prim = fitz.Rect(item[1])
                if (dtype == 'f'
                        and prim.width >= page_w * 0.9
                        and prim.height >= page_h * 0.9):
                    prim = None
            elif op == 'qu':
                quad = item[1]
                xs = [quad.ul.x, quad.ur.x, quad.ll.x, quad.lr.x]
                ys = [quad.ul.y, quad.ur.y, quad.ll.y, quad.lr.y]
                prim = fitz.Rect(min(xs), min(ys), max(xs), max(ys))
            elif op in ('c', 'v', 'y'):
                pts = [p for p in item[1:] if hasattr(p, 'x') and hasattr(p, 'y')]
                if pts:
                    xs = [p.x for p in pts]
                    ys = [p.y for p in pts]
                    prim = fitz.Rect(min(xs), min(ys), max(xs), max(ys))
            if prim is None:
                continue
            pr = _in_body(prim)
            if pr is None:
                continue
            if pr.width < 0.5 and pr.height < 0.5:
                continue
            drawing_rects.append(pr)

    image_rects = []
    try:
        for info in page.get_images(full=True):
            xref = info[0]
            try:
                rects = page.get_image_rects(xref)
            except Exception:
                rects = []
            for r in rects:
                rr = _in_body(r)
                if rr is not None:
                    image_rects.append(rr)
    except Exception:
        pass

    visual_rects = drawing_rects + image_rects

    if not visual_rects:
        doc.close()
        return _rasterize_by_whitespace(pdf_path, page_num, out_dir,
                                        has_bullets, dpi)

    x0 = min(r.x0 for r in visual_rects)
    y0 = min(r.y0 for r in visual_rects)
    x1 = max(r.x1 for r in visual_rects)
    y1 = max(r.y1 for r in visual_rects)
    bbox = fitz.Rect(x0, y0, x1, y1)

    drawings_y1 = max((r.y1 for r in drawing_rects), default=bbox.y1)
    drawings_y0 = min((r.y0 for r in drawing_rects), default=bbox.y0)

    text_dict = page.get_text("dict")
    for block in text_dict.get("blocks", []):
        if block.get("type", 1) != 0:
            continue
        bx = fitz.Rect(block["bbox"])
        if not bx.intersects(body_zone):
            continue
        if bx.y1 <= title_h or bx.y0 >= footer_y:
            continue
        if bx.y1 < bbox.y0 - 3 or bx.y0 > bbox.y1 + 3:
            continue
        if bx.x1 < bbox.x0 - 15 or bx.x0 > bbox.x1 + 15:
            continue
        if has_table and has_bullets and bx.y0 > drawings_y1 + 2:
            continue
        if has_table and bx.y1 < drawings_y0 - 12:
            continue
        bbox |= bx

    if has_bullets and drawing_rects:
        if bbox.y1 > drawings_y1 + 8:
            bbox.y1 = drawings_y1 + 4
        if has_table:
            bbox.y1 = min(bbox.y1, drawings_y1 + 2)

    pad_x = 6
    pad_y = 4
    bbox = fitz.Rect(
        max(0, bbox.x0 - pad_x),
        max(title_h, bbox.y0 - pad_y),
        min(page_w, bbox.x1 + pad_x),
        min(footer_y, bbox.y1 + pad_y),
    )

    if bbox.width < 20 or bbox.height < 20:
        doc.close()
        return _rasterize_by_whitespace(pdf_path, page_num, out_dir,
                                        has_bullets, dpi)

    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, clip=bbox)
    out_path = os.path.join(out_dir, f"slide_content_{page_num + 1}.png")
    pix.save(out_path)
    doc.close()
    return out_path


def _rasterize_by_whitespace(pdf_path, page_num, out_dir, has_bullets, dpi):
    """Fallback: full-page raster + whitespace-band crop.

    Used when the PDF has no usable drawings/images to infer a tight bbox
    (e.g. pure math rendered as text).
    """
    from PIL import Image as PILImage
    import numpy as np

    doc = fitz.open(pdf_path)
    if page_num >= len(doc):
        doc.close()
        return None

    page = doc[page_num]
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat)

    raw_path = os.path.join(out_dir, f"_raw_{page_num + 1}.png")
    pix.save(raw_path)
    doc.close()

    img = PILImage.open(raw_path)
    w, h = img.size
    arr = np.array(img)
    is_white = np.all(arr > 240, axis=2)

    white_bands = []
    in_band = False
    band_start = 0
    for r in range(h):
        row_white = is_white[r].all()
        if row_white and not in_band:
            band_start = r
            in_band = True
        elif not row_white and in_band:
            white_bands.append((band_start, r, r - band_start))
            in_band = False
    if in_band:
        white_bands.append((band_start, h, h - band_start))

    title_end = int(h * 0.12)
    for s, e, bh in white_bands:
        if s > int(h * 0.08) and bh >= 8:
            title_end = e
            break

    footer_start = int(h * 0.90)
    for s, e, bh in reversed(white_bands):
        if e < int(h * 0.98) and bh >= 20:
            footer_start = s
            break

    body_top = title_end
    body_bot = footer_start

    if has_bullets:
        best_gap = None
        best_gap_h = 0
        search_lo = int(body_top + (body_bot - body_top) * 0.25)
        search_hi = int(body_top + (body_bot - body_top) * 0.85)
        for s, e, bh in white_bands:
            if s >= search_lo and e <= search_hi and bh > best_gap_h:
                best_gap = (s, e)
                best_gap_h = bh
        if best_gap and best_gap_h >= 15:
            body_bot = best_gap[0]

    cropped = img.crop((0, body_top, w, body_bot))
    c_arr = np.array(cropped)
    c_white = np.all(c_arr > 240, axis=2)
    non_white_rows = np.where(~c_white.all(axis=1))[0]
    non_white_cols = np.where(~c_white.all(axis=0))[0]

    if len(non_white_rows) == 0 or len(non_white_cols) == 0:
        cropped.save(raw_path)
        return raw_path

    pad = int(dpi * 0.06)
    r_top = max(0, non_white_rows[0] - pad)
    r_bot = min(cropped.height, non_white_rows[-1] + pad)
    c_left = max(0, non_white_cols[0] - pad)
    c_right = min(cropped.width, non_white_cols[-1] + pad)

    trimmed = cropped.crop((c_left, r_top, c_right, r_bot))
    out_path = os.path.join(out_dir, f"slide_content_{page_num + 1}.png")
    trimmed.save(out_path)

    os.remove(raw_path)
    return out_path


# =====================================================================
# MAIN CONVERTER
# =====================================================================

def build_pptx(tex_path, pdf_path=None, output_path=None):
    base_dir = os.path.dirname(os.path.abspath(tex_path))
    tex = read_tex(tex_path)
    info = extract_preamble(tex)
    graphicspath = info.get('graphicspath', '')
    logo_path = None
    logo_m = re.search(
        r'\\includegraphics(?:\[.*?\])?\{([^{}]*logo[^{}]*)\}', tex
    )
    if logo_m:
        logo_path = resolve_image(logo_m.group(1), graphicspath, base_dir)
    if not logo_path:
        for name in (
            "Assets/logo.jpg", "Assets/logo.png", "Assets/logo.jpeg",
            "logo.jpg", "logo.png", "logo.jpeg",
        ):
            p = os.path.join(base_dir, name)
            if os.path.exists(p):
                logo_path = p
                break

    frames = extract_frames(tex)
    total = len(frames)

    if not output_path:
        output_path = os.path.join(base_dir,
            os.path.splitext(os.path.basename(tex_path))[0] + '_editable.pptx')

    extracted_dir = os.path.join(base_dir, 'Extracted_figures')
    fallback_dir = extracted_dir
    os.makedirs(extracted_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    LEFT_W = CONTENT_W * COL_SPLIT
    RIGHT_W = CONTENT_W * (1 - COL_SPLIT) - COL_GAP
    RIGHT_L = CONTENT_L + LEFT_W + COL_GAP

    footer_center = info.get("date") or FOOTER_TEXT_FALLBACK

    for idx, frame in enumerate(frames):
        slide_num = idx + 1
        s = prs.slides.add_slide(blank)

        # --- TITLE PAGE ---
        if frame['type'] == 'titlepage':
            _build_title_page(s, info)
            add_footer(s, slide_num, total, logo_path, footer_center)
            set_notes(s, frame.get('note', ''))
            continue

        title = frame.get('title', '')

        # --- FALLBACK: crop just the diagram/table from the PDF ---
        if frame.get('needs_fallback') and pdf_path:
            if title:
                add_title_bar(s, title)

            items = frame.get('items', [])
            is_enum = frame.get('is_enumerate', False)

            raw_content = frame.get('raw', '')
            has_table = '\\begin{tabular}' in raw_content
            fb_img = rasterize_content_region(
                pdf_path, idx, fallback_dir,
                has_bullets=bool(items),
                has_table=has_table,
            )
            if fb_img:
                from PIL import Image as PILImage
                with PILImage.open(fb_img) as im:
                    fb_aspect = im.width / im.height

                if items:
                    img_zone_h = BODY_H * 0.50
                    bullet_top = BODY_TOP + img_zone_h + Inches(0.1)
                    bullet_h = BODY_BOTTOM - bullet_top
                else:
                    img_zone_h = BODY_H * 0.75
                    bullet_top = None
                    bullet_h = 0

                avail_w = CONTENT_W
                if fb_aspect > (avail_w / img_zone_h):
                    iw = avail_w; ih = iw / fb_aspect
                else:
                    ih = img_zone_h; iw = ih * fb_aspect
                ix = CONTENT_L + (avail_w - iw) / 2
                iy = BODY_TOP + (img_zone_h - ih) / 2

                s.shapes.add_picture(fb_img, int(ix), int(iy),
                                     width=int(iw), height=int(ih))

                if items and bullet_top is not None:
                    add_bullet_list(s, items, CONTENT_L, bullet_top,
                                    CONTENT_W, bullet_h, is_enum)

            add_footer(s, slide_num, total, logo_path, footer_center)
            set_notes(s, frame.get('note', ''))
            continue

        # --- REGULAR CONTENT SLIDE ---
        if title:
            add_title_bar(s, title)

        items = frame.get('items', [])
        images = frame.get('images', [])
        captions = frame.get('captions', [])
        has_cols = frame.get('has_columns', False)
        is_enum = frame.get('is_enumerate', False)
        quote = frame.get('quote', '')

        if has_cols and images:
            if items:
                add_bullet_list(s, items, CONTENT_L, BODY_TOP, LEFT_W, BODY_H, is_enum)

            resolved = [resolve_image(img, graphicspath, base_dir) for img in images]
            resolved = [ensure_in_extracted(r, extracted_dir) for r in resolved if r]

            if len(resolved) == 1:
                _place_single_image(s, resolved[0], RIGHT_L, BODY_TOP, RIGHT_W, BODY_H)
                cap_text = captions[-1] if captions else ''
                if cap_text:
                    _add_caption(s, clean_text(cap_text), RIGHT_L, RIGHT_W)
            elif len(resolved) == 2 and len(captions) >= 2:
                clean_caps = [clean_text(c) for c in captions[:2]]
                _place_image_pair_vertical(s, resolved, RIGHT_L, BODY_TOP, RIGHT_W, BODY_H,
                                           captions=clean_caps)
            elif len(resolved) == 2:
                _place_image_pair_vertical(s, resolved, RIGHT_L, BODY_TOP, RIGHT_W, BODY_H)
                cap_text = captions[-1] if captions else ''
                if cap_text:
                    _add_caption(s, clean_text(cap_text), RIGHT_L, RIGHT_W)
            elif len(resolved) >= 3:
                _place_image_pair_vertical(s, resolved[:2], RIGHT_L, BODY_TOP, RIGHT_W, BODY_H * 0.85)

        elif has_cols and not images and items:
            mid = len(items) // 2
            left_items = items[:mid]
            right_items = items[mid:]
            half_w = (CONTENT_W - COL_GAP) / 2
            add_bullet_list(s, left_items, CONTENT_L, BODY_TOP, half_w, BODY_H, is_enum)
            add_bullet_list(s, right_items, CONTENT_L + half_w + COL_GAP, BODY_TOP, half_w, BODY_H, is_enum)

        elif items:
            bullet_h = BODY_H
            if quote:
                bullet_h = BODY_H - Inches(0.8)
            add_bullet_list(s, items, CONTENT_L, BODY_TOP, CONTENT_W, bullet_h, is_enum)
            if quote:
                add_quote(s, quote, BODY_BOTTOM - Inches(0.6))

        elif images:
            resolved = [resolve_image(img, graphicspath, base_dir) for img in images]
            resolved = [ensure_in_extracted(r, extracted_dir) for r in resolved if r]
            if resolved:
                _place_single_image(s, resolved[0], CONTENT_L, BODY_TOP, CONTENT_W, BODY_H)
                cap_text = captions[-1] if captions else ''
                if cap_text:
                    _add_caption_full(s, clean_text(cap_text))

        elif not title and 'raw' in frame:
            _build_closing_slide(s, frame['raw'])

        add_footer(s, slide_num, total, logo_path, footer_center)
        set_notes(s, frame.get('note', ''))

    apply_font_to_all(prs, FONT_NAME)
    prs.save(output_path)
    print(f"Saved {total} slides to: {output_path}")


# =====================================================================
# SLIDE BUILDERS
# =====================================================================

def _build_title_page(slide, info):
    title_text = info.get('title', 'Untitled')
    lines = title_text.split('\n')

    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), SLIDE_W - Inches(2.4), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True; _zero_insets(tf)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line.strip()
        r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = BLACK

    author_text = info.get('author', '')
    author_lines = [l.strip() for l in author_text.split('\n') if l.strip()]

    tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), SLIDE_W - Inches(2.4), Inches(1.6))
    tf = tb2.text_frame; tf.word_wrap = True; _zero_insets(tf)
    for i, line in enumerate(author_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4) if i > 0 else Pt(0)
        r = p.add_run(); r.text = line
        if i == 0:
            r.font.size = Pt(28)
        else:
            r.font.size = Pt(16)
        r.font.color.rgb = BLACK

    inst = info.get('institute', '')
    inst_lines = [l.strip() for l in inst.split('\n') if l.strip()]
    tb3 = slide.shapes.add_textbox(Inches(1.2), Inches(5.5), SLIDE_W - Inches(2.4), Inches(0.8))
    tf = tb3.text_frame; tf.word_wrap = True; _zero_insets(tf)
    for i, line in enumerate(inst_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(18); r.font.color.rgb = BLACK


def _build_closing_slide(slide, raw):
    """Build a centered text slide (e.g. Thank You)."""
    text = clean_text(raw)
    lines = [l.strip() for l in text.split('\n') if l.strip()]

    tb = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5), SLIDE_W - Inches(3), SLIDE_H - Inches(3))
    tf = tb.text_frame; tf.word_wrap = True; _zero_insets(tf)

    for i, line in enumerate(lines):
        is_bold = line.startswith('**') and line.endswith('**')
        if is_bold:
            line = line[2:-2]

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        if i == 0:
            font_sz = Pt(36)
        elif is_bold:
            font_sz = Pt(22)
        else:
            font_sz = Pt(18)

        _add_runs_from_text(p, line, font_sz)
        if i == 0 or is_bold:
            for run in p.runs:
                run.font.bold = True


def _place_single_image(slide, img_path, left, top, max_w, max_h):
    from PIL import Image as PILImage
    with PILImage.open(img_path) as img:
        aspect = img.width / img.height
    if aspect > (max_w / max_h):
        w = max_w
        h = w / aspect
    else:
        h = max_h
        w = h * aspect
    cx = left + (max_w - w) / 2
    cy = top + (max_h - h) / 2
    slide.shapes.add_picture(img_path, int(cx), int(cy), width=int(w), height=int(h))


def _place_image_pair_vertical(slide, images, left, top, max_w, max_h, captions=None):
    """Stack two images vertically in the right column with optional per-image captions."""
    from PIL import Image as PILImage
    cap_h = Inches(0.22) if captions and len(captions) >= 2 else 0
    gap = Inches(0.1)
    each_h = (max_h - gap - cap_h * min(len(captions or []), 2)) / 2
    for i, img_path in enumerate(images[:2]):
        with PILImage.open(img_path) as img:
            aspect = img.width / img.height
        if aspect > (max_w / each_h):
            w = max_w; h = w / aspect
        else:
            h = each_h; w = h * aspect
        slot_top = top + i * (each_h + gap + cap_h)
        cx = left + (max_w - w) / 2
        cy = slot_top + (each_h - h) / 2
        slide.shapes.add_picture(img_path, int(cx), int(cy), width=int(w), height=int(h))
        if captions and i < len(captions) and captions[i]:
            cap_y = slot_top + each_h
            tb = slide.shapes.add_textbox(int(left), int(cap_y), int(max_w), int(cap_h))
            tf = tb.text_frame; _zero_insets(tf)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = captions[i]
            r.font.size = FONT_CAPTION; r.font.italic = True; r.font.color.rgb = GRAY


def _add_caption(slide, text, left, width):
    cap_y = BODY_BOTTOM - Inches(0.4)
    tb = slide.shapes.add_textbox(left, cap_y, width, Inches(0.35))
    tf = tb.text_frame; _zero_insets(tf)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = FONT_CAPTION; r.font.italic = True; r.font.color.rgb = GRAY


def _add_caption_full(slide, text):
    _add_caption(slide, text, CONTENT_L, CONTENT_W)


# =====================================================================
# CLI
# =====================================================================

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <presentation.tex> [presentation.pdf] [output.pptx]")
        sys.exit(1)

    tex_path = sys.argv[1]
    pdf_path = sys.argv[2] if len(sys.argv) > 2 and sys.argv[2].endswith('.pdf') else None
    out_path = None
    for arg in sys.argv[2:]:
        if arg.endswith('.pptx'):
            out_path = arg

    if not pdf_path:
        default_pdf = os.path.splitext(tex_path)[0] + '.pdf'
        if os.path.exists(default_pdf):
            pdf_path = default_pdf
            print(f"Auto-detected PDF: {pdf_path}")

    build_pptx(tex_path, pdf_path, out_path)
