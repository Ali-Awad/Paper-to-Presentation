#!/usr/bin/env python3
"""Put slide images from slides_export/ into a PowerPoint file with speaker notes."""
import os
import re
import glob
from pptx import Presentation
from pptx.util import Emu

slide_dir = "slides_export"
out_pptx = "presentation.pptx"
tex_file = "presentation.tex"


def extract_notes(tex_path):
    """Extract \\note{...} blocks from the .tex file in order."""
    with open(tex_path, "r") as f:
        content = f.read()
    notes = []
    for m in re.finditer(r"\\note\{", content):
        start = m.end()
        depth = 1
        i = start
        while i < len(content) and depth > 0:
            if content[i] == "{":
                depth += 1
            elif content[i] == "}":
                depth -= 1
            i += 1
        notes.append(content[start : i - 1].strip())
    return notes


paths = sorted(glob.glob(os.path.join(slide_dir, "slide_*.png")))
if not paths:
    print(f"Error: No slide_*.png found in {slide_dir}. Run export_slides.py first.")
    exit(1)

notes = []
if os.path.isfile(tex_file):
    notes = extract_notes(tex_file)

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)
blank_slide_layout = prs.slide_layouts[6]

for idx, path in enumerate(paths):
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(
        path,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    if idx < len(notes) and notes[idx]:
        slide.notes_slide.notes_text_frame.text = notes[idx]
    print(f"Added {os.path.basename(path)}" + (" + notes" if idx < len(notes) and notes[idx] else ""))

prs.save(out_pptx)
print(f"Saved {out_pptx} ({len(paths)} slides, {min(len(notes), len(paths))} with notes)")
